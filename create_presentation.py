#!/usr/bin/env python3
"""
Create a PowerPoint (.pptx) using python-pptx.
Run: python3 create_presentation.py
Output: presentation.pptx in the current directory.

Edit the SLIDES list below to use your own content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt


# Customize your slides here: list of (title, bullet_points)
SLIDES = [
    ("Title Slide", ["Your presentation title", "Subtitle or date"]),
    ("Agenda", ["Point one", "Point two", "Point three"]),
    ("Section 1", ["Key message", "Supporting detail", "Another detail"]),
    ("Section 2", ["Key message", "Supporting detail"]),
    ("Thank You", ["Questions?", "Contact info"]),
]


def create_presentation(output_path="presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for title, bullets in SLIDES:
        # Title + content layout (layout 6 is often "Title and Content")
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True

        # Bullets
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.space_after = Pt(12)

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    create_presentation()
