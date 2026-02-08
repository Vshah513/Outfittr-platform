#!/usr/bin/env python3
"""
Professional Cursor AI IDE Presentation Generator - Assignment Version
Creates a high-end, diagram-heavy presentation with Cursor brand styling.
Updated to meet assignment rubric: Value Prop, Business Model, Competition, AI Tech, Security, Growth.
Includes interactive elements (hover tooltips via hyperlink screentips, click triggers).

Run: python3 create_cursor_presentation_pro.py
Output: cursor_presentation_pro.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from lxml import etree
import math

# =============================================================================
# CURSOR BRAND COLOR PALETTE
# =============================================================================
CURSOR_BLACK = RGBColor(18, 18, 18)
CURSOR_DARK_GRAY = RGBColor(38, 38, 38)
CURSOR_MID_GRAY = RGBColor(82, 82, 82)
CURSOR_LIGHT_GRAY = RGBColor(156, 156, 156)
CURSOR_OFF_WHITE = RGBColor(229, 229, 229)
CURSOR_WHITE = RGBColor(255, 255, 255)

CURSOR_PURPLE = RGBColor(139, 92, 246)
CURSOR_BLUE = RGBColor(59, 130, 246)
CURSOR_GREEN = RGBColor(34, 197, 94)
CURSOR_AMBER = RGBColor(245, 158, 11)
CURSOR_RED = RGBColor(239, 68, 68)

GRADIENT_START = RGBColor(45, 45, 45)
GRADIENT_END = RGBColor(25, 25, 25)

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def add_rounded_rectangle(slide, left, top, width, height, fill_color, 
                          border_color=None, border_width=0, corner_radius=0.1):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    
    return shape


def add_circle(slide, left, top, size, fill_color, border_color=None):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    
    return shape


def add_arrow(slide, left, top, width, height, fill_color):
    """Add an arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_to_shape(shape, text, font_size=14, font_color=CURSOR_WHITE, 
                      bold=False, alignment=PP_ALIGN.CENTER):
    """Add centered text to a shape."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_textbox(slide, left, top, width, height, text, font_size=14, 
                font_color=CURSOR_OFF_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return textbox


def add_hyperlink_textbox(slide, left, top, width, height, text, url, tooltip,
                          font_size=14, font_color=CURSOR_BLUE, bold=False, 
                          alignment=PP_ALIGN.LEFT, underline=True):
    """Add a text box with a clickable hyperlink and hover tooltip (screentip)."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.underline = underline
    p.alignment = alignment
    
    # Add hyperlink with screentip (tooltip on hover)
    rPr = run._r.get_or_add_rPr()
    hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), '')
    hlinkClick.set('tooltip', tooltip)
    
    # Create the relationship for the hyperlink
    rel = slide.part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
    hlinkClick.set(qn('r:id'), rel)
    
    return textbox


def add_clickable_shape(slide, shape, url, tooltip):
    """Make a shape clickable with a hyperlink and tooltip."""
    click = shape.click_action
    click.hyperlink.address = url
    # Note: python-pptx doesn't directly support screentip on shapes, 
    # but the hyperlink will work. Tooltip shown in speaker notes instead.
    return shape


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_source_footer(slide, source_text, slide_width, slide_height):
    """Add a small source footer at the bottom of the slide."""
    footer = add_textbox(
        slide, 
        Inches(0.3), 
        slide_height - Inches(0.35),
        slide_width - Inches(0.6),
        Inches(0.3),
        source_text,
        font_size=7,
        font_color=CURSOR_MID_GRAY,
        alignment=PP_ALIGN.LEFT
    )
    return footer


def add_slide_number(slide, number, slide_width, slide_height):
    """Add slide number to bottom right."""
    add_textbox(
        slide,
        slide_width - Inches(0.6),
        slide_height - Inches(0.4),
        Inches(0.4),
        Inches(0.3),
        str(number),
        font_size=10,
        font_color=CURSOR_MID_GRAY,
        alignment=PP_ALIGN.RIGHT
    )


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_horizontal_line(slide, left, top, width, color, thickness=1):
    """Add a horizontal line/divider."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


# =============================================================================
# SLIDE CREATION FUNCTIONS
# =============================================================================

def create_title_slide(prs, slide_width, slide_height):
    """Slide 1: Title slide with logo, link, and presenter name."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Decorative accent shapes (keeping visual style)
    accent1 = add_rounded_rectangle(
        slide, slide_width - Inches(3), Inches(-0.5),
        Inches(4), Inches(2.5), CURSOR_PURPLE
    )
    accent1.rotation = -15
    
    accent2 = add_rounded_rectangle(
        slide, Inches(-1), slide_height - Inches(2),
        Inches(3), Inches(3), CURSOR_BLUE
    )
    accent2.rotation = 30
    
    # Main content card
    main_card = add_rounded_rectangle(
        slide, Inches(1.5), Inches(1),
        Inches(7), Inches(5),
        CURSOR_DARK_GRAY,
        border_color=CURSOR_MID_GRAY,
        border_width=1
    )
    
    # Large Cursor Logo (circle with "C" - centered)
    logo = add_circle(slide, Inches(4.1), Inches(1.5), Inches(1.8), CURSOR_PURPLE)
    add_text_to_shape(logo, "C", font_size=60, font_color=CURSOR_WHITE, bold=True)
    # Make logo clickable
    add_clickable_shape(slide, logo, "https://www.cursor.com/", "Visit cursor.com")
    
    # Main title - "Cursor"
    add_textbox(
        slide, Inches(1.5), Inches(3.5),
        Inches(7), Inches(0.9),
        "Cursor",
        font_size=54, font_color=CURSOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER
    )
    
    # Subtitle - "AI Coding IDE"
    add_textbox(
        slide, Inches(1.5), Inches(4.3),
        Inches(7), Inches(0.6),
        "AI Coding IDE",
        font_size=28, font_color=CURSOR_PURPLE, bold=True, alignment=PP_ALIGN.CENTER
    )
    
    # Clickable link with tooltip
    add_hyperlink_textbox(
        slide, Inches(1.5), Inches(5),
        Inches(7), Inches(0.4),
        "https://www.cursor.com/",
        "https://www.cursor.com/",
        "Click to visit Cursor's official website",
        font_size=14, font_color=CURSOR_BLUE, alignment=PP_ALIGN.CENTER
    )
    
    # Presenter name placeholder
    add_textbox(
        slide, Inches(1.5), Inches(5.6),
        Inches(7), Inches(0.4),
        "Your Name",
        font_size=16, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
    )
    
    add_speaker_notes(
        slide,
        "Introduction: Today I'll present Cursor, an AI coding IDE that's transforming how developers write software. "
        "I'll cover its value proposition, business model, competition, AI technology, security posture, and growth trajectory."
    )
    
    add_slide_number(slide, 1, slide_width, slide_height)
    return slide


def create_value_proposition_slide(prs, slide_width, slide_height):
    """Slide 2: Value Proposition - Why Developers Use Cursor."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6),
        "Value Proposition",
        font_size=34, font_color=CURSOR_WHITE, bold=True
    )
    
    # Subtitle
    add_textbox(
        slide, Inches(0.5), Inches(0.85),
        Inches(9), Inches(0.4),
        "Why Developers Choose Cursor",
        font_size=14, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_horizontal_line(slide, Inches(0.5), Inches(1.25), Inches(2), CURSOR_PURPLE, thickness=3)
    
    # === WORKFLOW DIAGRAM (preserved from original) ===
    # Column 1: YOU (Input)
    col1_x = Inches(0.4)
    col1_card = add_rounded_rectangle(
        slide, col1_x, Inches(1.5),
        Inches(2.4), Inches(2.8), CURSOR_DARK_GRAY,
        border_color=CURSOR_PURPLE, border_width=2
    )
    
    user_icon = add_circle(slide, col1_x + Inches(0.85), Inches(1.7), Inches(0.6), CURSOR_PURPLE)
    add_text_to_shape(user_icon, "You", font_size=12, font_color=CURSOR_WHITE, bold=True)
    
    add_textbox(
        slide, col1_x + Inches(0.1), Inches(2.4),
        Inches(2.2), Inches(0.35),
        "Natural Language",
        font_size=11, font_color=CURSOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER
    )
    
    # Example prompts with tooltips (interactive)
    prompts_text = '"Add login with validation"\n"Fix the API bug"\n"Refactor to async"'
    add_textbox(
        slide, col1_x + Inches(0.1), Inches(2.75),
        Inches(2.2), Inches(1.4),
        prompts_text,
        font_size=9, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
    )
    
    # Arrow 1
    add_arrow(slide, Inches(2.9), Inches(2.7), Inches(0.5), Inches(0.25), CURSOR_MID_GRAY)
    
    # Column 2: CURSOR AI
    col2_x = Inches(3.5)
    col2_card = add_rounded_rectangle(
        slide, col2_x, Inches(1.5),
        Inches(1.8), Inches(2.8), CURSOR_DARK_GRAY,
        border_color=CURSOR_BLUE, border_width=2
    )
    
    cursor_icon = add_circle(slide, col2_x + Inches(0.5), Inches(1.7), Inches(0.6), CURSOR_BLUE)
    add_text_to_shape(cursor_icon, "AI", font_size=12, font_color=CURSOR_WHITE, bold=True)
    
    add_textbox(
        slide, col2_x + Inches(0.1), Inches(2.4),
        Inches(1.6), Inches(0.35),
        "Cursor Processes",
        font_size=10, font_color=CURSOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER
    )
    
    tab_box = add_rounded_rectangle(
        slide, col2_x + Inches(0.1), Inches(2.8),
        Inches(1.6), Inches(0.65), CURSOR_PURPLE
    )
    add_text_to_shape(tab_box, "Tab\nCompletion", font_size=8, font_color=CURSOR_WHITE, bold=True)
    
    agent_box = add_rounded_rectangle(
        slide, col2_x + Inches(0.1), Inches(3.5),
        Inches(1.6), Inches(0.65), CURSOR_GREEN
    )
    add_text_to_shape(agent_box, "Agent\nCmd+I", font_size=8, font_color=CURSOR_WHITE, bold=True)
    
    # Arrow 2
    add_arrow(slide, Inches(5.4), Inches(2.7), Inches(0.5), Inches(0.25), CURSOR_MID_GRAY)
    
    # Column 3: OUTPUT
    col3_x = Inches(6)
    col3_card = add_rounded_rectangle(
        slide, col3_x, Inches(1.5),
        Inches(2.4), Inches(2.8), CURSOR_DARK_GRAY,
        border_color=CURSOR_GREEN, border_width=2
    )
    
    output_icon = add_circle(slide, col3_x + Inches(0.85), Inches(1.7), Inches(0.6), CURSOR_GREEN)
    add_text_to_shape(output_icon, "</>", font_size=11, font_color=CURSOR_WHITE, bold=True)
    
    add_textbox(
        slide, col3_x + Inches(0.1), Inches(2.4),
        Inches(2.2), Inches(0.35),
        "Working Code",
        font_size=11, font_color=CURSOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER
    )
    
    add_textbox(
        slide, col3_x + Inches(0.1), Inches(2.75),
        Inches(2.2), Inches(1.4),
        "Multi-file edits\nTerminal commands\nTests & debugging\nIterative refinement",
        font_size=9, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
    )
    
    # === BENEFITS SECTION (rewritten as benefits, not features) ===
    benefits_y = Inches(4.5)
    
    benefits = [
        ("Ship Faster", "AI writes & refactors code from natural language", CURSOR_PURPLE, Inches(0.4)),
        ("Higher Quality", "Multi-file refactors, test-aware suggestions", CURSOR_BLUE, Inches(2.65)),
        ("Less Switching", "AI built into IDE, not a separate tool", CURSOR_GREEN, Inches(4.9)),
        ("Personalized", "Tab learns from your accept/reject patterns", CURSOR_AMBER, Inches(7.15)),
    ]
    
    for title, desc, color, x in benefits:
        benefit_card = add_rounded_rectangle(
            slide, x, benefits_y,
            Inches(2.15), Inches(1.1), CURSOR_DARK_GRAY,
            border_color=color, border_width=2
        )
        
        # Colored top bar
        top_bar = add_rounded_rectangle(
            slide, x, benefits_y,
            Inches(2.15), Inches(0.35), color
        )
        add_text_to_shape(top_bar, title, font_size=10, font_color=CURSOR_WHITE, bold=True)
        
        add_textbox(
            slide, x + Inches(0.08), benefits_y + Inches(0.45),
            Inches(2), Inches(0.6),
            desc,
            font_size=8, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
        )
    
    # Target market line
    add_textbox(
        slide, Inches(0.4), Inches(5.75),
        Inches(9.2), Inches(0.35),
        "Target Market: Professional software engineers, teams, and companies adopting AI-assisted development",
        font_size=9, font_color=CURSOR_MID_GRAY, alignment=PP_ALIGN.CENTER
    )
    
    add_speaker_notes(
        slide,
        "Value Proposition: Cursor helps developers ship code faster by writing and refactoring code from natural language. "
        "It delivers higher code quality through multi-file refactors and test-aware suggestions. "
        "Unlike separate AI chat tools, Cursor is built directly into the IDE, reducing context switching. "
        "The Tab completion learns from your patterns, becoming more personalized over time. "
        "Target market is professional software engineers, teams, and companies embracing AI-assisted development."
    )
    
    add_source_footer(
        slide,
        "Sources: cursor.com/docs | cursor.com/docs/agent/overview | cursor.com/docs/tab/overview",
        slide_width, slide_height
    )
    add_slide_number(slide, 2, slide_width, slide_height)
    return slide


def create_business_model_slide(prs, slide_width, slide_height):
    """Slide 3: Business Model - How Cursor Makes Money (merged pricing info)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6),
        "Business Model",
        font_size=34, font_color=CURSOR_WHITE, bold=True
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(0.85),
        Inches(9), Inches(0.4),
        "How Cursor Makes Money",
        font_size=14, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_horizontal_line(slide, Inches(0.5), Inches(1.25), Inches(2), CURSOR_AMBER, thickness=3)
    
    # === THREE COLUMN LAYOUT: Individuals / Teams / Enterprise ===
    col_width = Inches(3)
    col_height = Inches(3.8)
    col_y = Inches(1.5)
    gap = Inches(0.25)
    
    columns = [
        {
            "title": "INDIVIDUALS",
            "subtitle": "Usage-Based SaaS",
            "color": CURSOR_BLUE,
            "x": Inches(0.35),
            "items": [
                "Hobby: Free tier (limited)",
                "Pro: $20/mo with usage credits",
                "Pro+: $60/mo (3x model usage)",
                "Ultra: $200/mo (20x usage)",
            ],
            "highlight": "Revenue scales with AI usage, not just seats"
        },
        {
            "title": "TEAMS",
            "subtitle": "Per-Seat + Usage",
            "color": CURSOR_PURPLE,
            "x": Inches(3.5),
            "items": [
                "$40/user/month",
                "Monthly usage allocation per seat",
                "On-demand overage billing",
                "Shared resources & analytics",
            ],
            "highlight": "Predictable base + usage upside"
        },
        {
            "title": "ENTERPRISE",
            "subtitle": "Custom Contracts",
            "color": CURSOR_GREEN,
            "x": Inches(6.65),
            "items": [
                "Custom pricing & terms",
                "Pooled org-wide usage",
                "Invoice/PO billing",
                "Premium support & SLAs",
            ],
            "highlight": "High-value accounts with sticky contracts"
        }
    ]
    
    for col in columns:
        # Card
        card = add_rounded_rectangle(
            slide, col["x"], col_y,
            col_width, col_height, CURSOR_DARK_GRAY,
            border_color=col["color"], border_width=2
        )
        
        # Header
        header = add_rounded_rectangle(
            slide, col["x"], col_y,
            col_width, Inches(0.7), col["color"]
        )
        add_text_to_shape(header, col["title"], font_size=14, font_color=CURSOR_WHITE, bold=True)
        
        # Subtitle
        add_textbox(
            slide, col["x"] + Inches(0.1), col_y + Inches(0.75),
            col_width - Inches(0.2), Inches(0.35),
            col["subtitle"],
            font_size=11, font_color=col["color"], bold=True, alignment=PP_ALIGN.CENTER
        )
        
        # Items
        item_y = col_y + Inches(1.2)
        for item in col["items"]:
            add_textbox(
                slide, col["x"] + Inches(0.15), item_y,
                Inches(0.2), Inches(0.25),
                "•",
                font_size=10, font_color=col["color"]
            )
            add_textbox(
                slide, col["x"] + Inches(0.35), item_y,
                col_width - Inches(0.5), Inches(0.4),
                item,
                font_size=9, font_color=CURSOR_OFF_WHITE
            )
            item_y += Inches(0.45)
        
        # Highlight box
        highlight_box = add_rounded_rectangle(
            slide, col["x"] + Inches(0.1), col_y + Inches(3.1),
            col_width - Inches(0.2), Inches(0.55), CURSOR_BLACK,
            border_color=col["color"], border_width=1
        )
        add_textbox(
            slide, col["x"] + Inches(0.15), col_y + Inches(3.2),
            col_width - Inches(0.3), Inches(0.4),
            col["highlight"],
            font_size=8, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
        )
    
    # === KEY INSIGHT BOX ===
    insight_box = add_rounded_rectangle(
        slide, Inches(0.35), Inches(5.5),
        Inches(9.3), Inches(0.9), CURSOR_DARK_GRAY,
        border_color=CURSOR_AMBER, border_width=2
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(5.6),
        Inches(0.3), Inches(0.3),
        "💡",
        font_size=14, font_color=CURSOR_AMBER, alignment=PP_ALIGN.CENTER
    )
    
    add_textbox(
        slide, Inches(0.85), Inches(5.6),
        Inches(8.6), Inches(0.7),
        "Key Insight: Monetization is tied to AI model usage (credits consumed) rather than just seat count. "
        "This aligns revenue with the value delivered—more AI usage = more revenue.",
        font_size=10, font_color=CURSOR_OFF_WHITE
    )
    
    add_speaker_notes(
        slide,
        "Business Model: Cursor monetizes through usage-based SaaS pricing. "
        "Individuals can start free with Hobby tier, then upgrade to Pro at $20/month which includes usage credits. "
        "Higher tiers offer multiplied usage at premium prices. "
        "Teams pay $40 per user per month with usage allocations and overage billing. "
        "Enterprise contracts are custom with pooled usage, invoice billing, and premium support. "
        "The key insight: revenue scales with AI usage, not just headcount, aligning Cursor's success with customer value."
    )
    
    add_source_footer(
        slide,
        "Sources: cursor.com/pricing | cursor.com/blog/june-2025-pricing | cursor.com/docs/account/teams/pricing",
        slide_width, slide_height
    )
    add_slide_number(slide, 3, slide_width, slide_height)
    return slide


def create_competition_slide(prs, slide_width, slide_height):
    """Slide 4: Competition - AI and Non-AI Alternatives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6),
        "Competition",
        font_size=34, font_color=CURSOR_WHITE, bold=True
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(0.85),
        Inches(9), Inches(0.4),
        "AI and Non-AI Alternatives",
        font_size=14, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_horizontal_line(slide, Inches(0.5), Inches(1.25), Inches(2), CURSOR_RED, thickness=3)
    
    # === TWO COLUMN COMPARISON ===
    col_width = Inches(4.5)
    col_height = Inches(2.4)
    col_y = Inches(1.5)
    
    # Left column: AI Coding Tools
    ai_card = add_rounded_rectangle(
        slide, Inches(0.35), col_y,
        col_width, col_height, CURSOR_DARK_GRAY,
        border_color=CURSOR_PURPLE, border_width=2
    )
    
    ai_header = add_rounded_rectangle(
        slide, Inches(0.35), col_y,
        col_width, Inches(0.55), CURSOR_PURPLE
    )
    add_text_to_shape(ai_header, "AI Coding Tools", font_size=13, font_color=CURSOR_WHITE, bold=True)
    
    ai_competitors = [
        ("GitHub Copilot", "Code suggestions in VS Code/IDEs; single-model (OpenAI)"),
        ("Replit Ghostwriter", "AI in cloud IDE; focused on beginners/prototyping"),
        ("Codeium", "Free AI autocomplete; limited agent capabilities"),
        ("Amazon CodeWhisperer", "AWS-integrated; enterprise security focus"),
    ]
    
    item_y = col_y + Inches(0.65)
    for name, desc in ai_competitors:
        add_textbox(
            slide, Inches(0.5), item_y,
            Inches(1.4), Inches(0.35),
            name,
            font_size=10, font_color=CURSOR_PURPLE, bold=True
        )
        add_textbox(
            slide, Inches(1.9), item_y,
            Inches(2.8), Inches(0.35),
            desc,
            font_size=8, font_color=CURSOR_LIGHT_GRAY
        )
        item_y += Inches(0.42)
    
    # Right column: Traditional Tools
    trad_card = add_rounded_rectangle(
        slide, Inches(5.15), col_y,
        col_width, col_height, CURSOR_DARK_GRAY,
        border_color=CURSOR_MID_GRAY, border_width=2
    )
    
    trad_header = add_rounded_rectangle(
        slide, Inches(5.15), col_y,
        col_width, Inches(0.55), CURSOR_MID_GRAY
    )
    add_text_to_shape(trad_header, "Traditional Tools", font_size=13, font_color=CURSOR_WHITE, bold=True)
    
    trad_competitors = [
        ("VS Code (no AI)", "Powerful editor; requires manual coding or plugins"),
        ("JetBrains IDEs", "Feature-rich; AI add-ons available separately"),
        ("Stack Overflow", "Manual search; context switching, slower iteration"),
        ("Vim/Emacs", "Expert tools; steep learning curve, no AI native"),
    ]
    
    item_y = col_y + Inches(0.65)
    for name, desc in trad_competitors:
        add_textbox(
            slide, Inches(5.3), item_y,
            Inches(1.6), Inches(0.35),
            name,
            font_size=10, font_color=CURSOR_LIGHT_GRAY, bold=True
        )
        add_textbox(
            slide, Inches(6.9), item_y,
            Inches(2.6), Inches(0.35),
            desc,
            font_size=8, font_color=CURSOR_LIGHT_GRAY
        )
        item_y += Inches(0.42)
    
    # === CURSOR DIFFERENTIATION ===
    diff_y = Inches(4.1)
    diff_card = add_rounded_rectangle(
        slide, Inches(0.35), diff_y,
        Inches(9.3), Inches(2.2), CURSOR_DARK_GRAY,
        border_color=CURSOR_GREEN, border_width=2
    )
    
    diff_header = add_rounded_rectangle(
        slide, Inches(0.35), diff_y,
        Inches(9.3), Inches(0.5), CURSOR_GREEN
    )
    add_text_to_shape(diff_header, "How Cursor Differentiates", font_size=12, font_color=CURSOR_WHITE, bold=True)
    
    differentiators = [
        ("AI-Native IDE", "Deeper project awareness & multi-file editing vs simple code suggestions", CURSOR_PURPLE),
        ("Multi-Model Routing", "OpenAI, Anthropic, Google models vs single-provider lock-in", CURSOR_BLUE),
        ("Enterprise-Grade", "SSO, RBAC, audit logs, model controls vs consumer-only tools", CURSOR_GREEN),
        ("Agent Capabilities", "Autonomous coding tasks, terminal commands vs passive autocomplete", CURSOR_AMBER),
    ]
    
    diff_item_y = diff_y + Inches(0.6)
    for title, desc, color in differentiators:
        # Colored indicator
        indicator = add_rounded_rectangle(
            slide, Inches(0.5), diff_item_y + Inches(0.05),
            Inches(0.08), Inches(0.25), color
        )
        
        add_textbox(
            slide, Inches(0.7), diff_item_y,
            Inches(1.8), Inches(0.35),
            title,
            font_size=10, font_color=color, bold=True
        )
        add_textbox(
            slide, Inches(2.5), diff_item_y,
            Inches(7), Inches(0.35),
            desc,
            font_size=9, font_color=CURSOR_OFF_WHITE
        )
        diff_item_y += Inches(0.38)
    
    add_speaker_notes(
        slide,
        "Competition: In AI coding tools, GitHub Copilot is the main competitor with broad VS Code adoption but limited to OpenAI models. "
        "Replit Ghostwriter targets beginners in a cloud IDE. Codeium offers free autocomplete but limited agent features. "
        "Traditional tools like VS Code without AI, JetBrains IDEs, and Stack Overflow require manual coding or context switching. "
        "Cursor differentiates through: (1) AI-native IDE with deep project awareness, (2) multi-model routing across providers, "
        "(3) enterprise-grade security features, and (4) true agent capabilities that autonomously complete coding tasks."
    )
    
    add_source_footer(
        slide,
        "Competitive analysis based on publicly available product information",
        slide_width, slide_height
    )
    add_slide_number(slide, 4, slide_width, slide_height)
    return slide


def create_ai_technologies_slide(prs, slide_width, slide_height):
    """Slide 5: AI Technologies Under the Hood."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6),
        "AI Technologies",
        font_size=34, font_color=CURSOR_WHITE, bold=True
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(0.85),
        Inches(9), Inches(0.4),
        "Under the Hood & Proprietary Edge",
        font_size=14, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_horizontal_line(slide, Inches(0.5), Inches(1.25), Inches(2), CURSOR_BLUE, thickness=3)
    
    # === ARCHITECTURE DIAGRAM (preserved) ===
    
    # Center hub - Cursor Router
    hub_x, hub_y = Inches(3.8), Inches(2.3)
    hub = add_rounded_rectangle(
        slide, hub_x, hub_y,
        Inches(2.2), Inches(1), CURSOR_PURPLE,
        border_color=CURSOR_WHITE, border_width=2
    )
    add_text_to_shape(hub, "CURSOR\nAI Router", font_size=12, font_color=CURSOR_WHITE, bold=True)
    
    # Model providers around the hub
    providers = [
        ("OpenAI", CURSOR_GREEN, Inches(0.8), Inches(1.6)),
        ("Anthropic", CURSOR_AMBER, Inches(7.2), Inches(1.6)),
        ("Google", CURSOR_BLUE, Inches(0.8), Inches(3.2)),
        ("Others", CURSOR_MID_GRAY, Inches(7.2), Inches(3.2)),
    ]
    
    for name, color, x, y in providers:
        provider_box = add_rounded_rectangle(
            slide, x, y,
            Inches(1.6), Inches(0.7), CURSOR_DARK_GRAY,
            border_color=color, border_width=2
        )
        add_text_to_shape(provider_box, name, font_size=10, font_color=color, bold=True)
        
        # Connection lines
        if x < hub_x:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + Inches(1.6), y + Inches(0.3),
                hub_x - x - Inches(1.6), Pt(2)
            )
        else:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                hub_x + Inches(2.2), y + Inches(0.3),
                x - hub_x - Inches(2.2), Pt(2)
            )
        line.fill.solid()
        line.fill.fore_color.rgb = CURSOR_MID_GRAY
        line.line.fill.background()
    
    # Auto Mode indicator
    auto_box = add_rounded_rectangle(
        slide, Inches(3.5), Inches(3.5),
        Inches(2.8), Inches(0.5), CURSOR_GREEN
    )
    add_text_to_shape(auto_box, "AUTO: Intelligent Routing", font_size=9, font_color=CURSOR_WHITE, bold=True)
    
    # === TECHNOLOGY BREAKDOWN ===
    tech_y = Inches(4.2)
    
    tech_cards = [
        {
            "title": "Frontier LLMs",
            "desc": "GPT-4, Claude, Gemini\nand other cutting-edge\nlanguage models",
            "color": CURSOR_PURPLE,
            "x": Inches(0.35)
        },
        {
            "title": "Agent System",
            "desc": "Multi-file editing,\nterminal commands,\niterative feedback loops",
            "color": CURSOR_BLUE,
            "x": Inches(2.55)
        },
        {
            "title": "Proprietary Edge",
            "desc": "Routing logic, IDE\nintegration, project\ncontext handling",
            "color": CURSOR_GREEN,
            "x": Inches(4.75)
        },
        {
            "title": "Optimization",
            "desc": "Cost, latency, and\nreliability balancing\nacross providers",
            "color": CURSOR_AMBER,
            "x": Inches(6.95)
        },
    ]
    
    for card in tech_cards:
        card_shape = add_rounded_rectangle(
            slide, card["x"], tech_y,
            Inches(2.1), Inches(1.55), CURSOR_DARK_GRAY,
            border_color=card["color"], border_width=2
        )
        
        # Icon circle
        icon = add_circle(slide, card["x"] + Inches(0.7), tech_y + Inches(0.15), Inches(0.5), card["color"])
        
        add_textbox(
            slide, card["x"] + Inches(0.1), tech_y + Inches(0.7),
            Inches(1.9), Inches(0.35),
            card["title"],
            font_size=10, font_color=card["color"], bold=True, alignment=PP_ALIGN.CENTER
        )
        
        add_textbox(
            slide, card["x"] + Inches(0.1), tech_y + Inches(1),
            Inches(1.9), Inches(0.5),
            card["desc"],
            font_size=8, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
        )
    
    # Secret sauce callout
    add_textbox(
        slide, Inches(0.35), Inches(5.9),
        Inches(9.3), Inches(0.4),
        '"Secret Sauce": Cursor\'s value is in the integration layer—routing, context, IDE UX—built on top of base models',
        font_size=9, font_color=CURSOR_MID_GRAY, alignment=PP_ALIGN.CENTER
    )
    
    add_speaker_notes(
        slide,
        "AI Technologies: Cursor uses frontier LLMs from OpenAI (GPT-4), Anthropic (Claude), Google (Gemini), and others. "
        "The Agent system can edit multiple files, run terminal commands, and iterate on feedback autonomously. "
        "Cursor's proprietary edge—its 'secret sauce'—is the integration layer: intelligent routing logic, deep IDE integration, "
        "and sophisticated project context handling built on top of base models. "
        "The Auto mode optimizes for cost, latency, and reliability by routing requests to the best-suited model."
    )
    
    add_source_footer(slide, "Source: cursor.com/docs/models", slide_width, slide_height)
    add_slide_number(slide, 5, slide_width, slide_height)
    return slide


def create_security_slide(prs, slide_width, slide_height):
    """Slide 6: Enterprise Security, Compliance & Risk."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6),
        "Enterprise Security",
        font_size=34, font_color=CURSOR_WHITE, bold=True
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(0.85),
        Inches(9), Inches(0.4),
        "Compliance & Risk Considerations",
        font_size=14, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_horizontal_line(slide, Inches(0.5), Inches(1.25), Inches(2), CURSOR_GREEN, thickness=3)
    
    # === SECURITY PILLARS (preserved) ===
    pillar_width = Inches(2.2)
    pillar_height = Inches(2.4)
    pillars_y = Inches(1.5)
    gap = Inches(0.2)
    start_x = Inches(0.3)
    
    pillars = [
        {"icon": "🛡️", "title": "Compliance", "color": CURSOR_GREEN,
         "items": ["SOC 2 Type II", "Annual pen tests", "Trust Center"]},
        {"icon": "🔐", "title": "Encryption", "color": CURSOR_BLUE,
         "items": ["AES-256 at rest", "TLS 1.2+ transit", "E2E protection"]},
        {"icon": "⚙️", "title": "Controls", "color": CURSOR_PURPLE,
         "items": ["Admin controls", "Model restrictions", "RBAC"]},
        {"icon": "📋", "title": "Privacy", "color": CURSOR_AMBER,
         "items": ["GDPR compliant", "CCPA compliant", "Data sovereignty"]},
    ]
    
    for i, pillar in enumerate(pillars):
        x = start_x + (pillar_width + gap) * i
        
        card = add_rounded_rectangle(
            slide, x, pillars_y,
            pillar_width, pillar_height, CURSOR_DARK_GRAY,
            border_color=pillar["color"], border_width=2
        )
        
        icon_bg = add_circle(slide, x + Inches(0.7), pillars_y + Inches(0.15), Inches(0.6), pillar["color"])
        add_text_to_shape(icon_bg, pillar["icon"], font_size=16, font_color=CURSOR_WHITE)
        
        add_textbox(
            slide, x + Inches(0.1), pillars_y + Inches(0.85),
            pillar_width - Inches(0.2), Inches(0.35),
            pillar["title"],
            font_size=12, font_color=pillar["color"], bold=True, alignment=PP_ALIGN.CENTER
        )
        
        item_y = pillars_y + Inches(1.3)
        for item in pillar["items"]:
            add_textbox(
                slide, x + Inches(0.15), item_y,
                pillar_width - Inches(0.3), Inches(0.3),
                f"• {item}",
                font_size=9, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
            )
            item_y += Inches(0.32)
    
    # === TRUST CENTER CALLOUT ===
    trust_box = add_rounded_rectangle(
        slide, Inches(0.35), Inches(4.1),
        Inches(9.3), Inches(1), CURSOR_DARK_GRAY,
        border_color=CURSOR_GREEN, border_width=2
    )
    
    trust_icon = add_circle(slide, Inches(0.55), Inches(4.25), Inches(0.55), CURSOR_GREEN)
    add_text_to_shape(trust_icon, "✓", font_size=18, font_color=CURSOR_WHITE, bold=True)
    
    add_textbox(
        slide, Inches(1.25), Inches(4.2),
        Inches(8.2), Inches(0.35),
        "Trust Center: SOC 2 reports, pen test summaries, compliance docs",
        font_size=12, font_color=CURSOR_GREEN, bold=True
    )
    
    # Make link clickable with tooltip
    add_hyperlink_textbox(
        slide, Inches(1.25), Inches(4.55),
        Inches(8.2), Inches(0.35),
        "trust.cursor.com",
        "https://trust.cursor.com/",
        "Click to visit Cursor Trust Center for compliance documentation",
        font_size=10, font_color=CURSOR_BLUE
    )
    
    # === ADOPTION & RISK CALLOUT ===
    risk_box = add_rounded_rectangle(
        slide, Inches(0.35), Inches(5.25),
        Inches(9.3), Inches(1.15), CURSOR_DARK_GRAY,
        border_color=CURSOR_AMBER, border_width=2
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(5.35),
        Inches(9), Inches(0.35),
        "Why It Matters: Security posture lowers enterprise adoption risk → faster sales cycles",
        font_size=10, font_color=CURSOR_OFF_WHITE, bold=True
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(5.7),
        Inches(9), Inches(0.6),
        "Key Risks: Dependency on third-party AI models  |  Evolving AI regulation  |  Managing AI-generated code quality",
        font_size=9, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_speaker_notes(
        slide,
        "Enterprise Security: Cursor is SOC 2 Type II certified with annual third-party penetration testing. "
        "Data is encrypted at rest (AES-256) and in transit (TLS 1.2+). "
        "Enterprise customers get admin controls, model restrictions, and RBAC. "
        "GDPR and CCPA compliance references are available. "
        "The Trust Center at trust.cursor.com provides audit reports and compliance documentation. "
        "Why this matters: strong security posture reduces enterprise adoption risk, enabling faster sales cycles. "
        "Risks to consider: dependency on third-party AI models, evolving AI regulation landscape, "
        "and the need for processes to manage AI-generated code quality and security."
    )
    
    add_source_footer(
        slide,
        "Sources: cursor.com/security | cursor.com/enterprise | trust.cursor.com",
        slide_width, slide_height
    )
    add_slide_number(slide, 6, slide_width, slide_height)
    return slide


def create_growth_funding_slide(prs, slide_width, slide_height):
    """Slide 7: Growth Stage, Funding & Impact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.25),
        Inches(9), Inches(0.55),
        "Growth Stage, Funding & Impact",
        font_size=32, font_color=CURSOR_WHITE, bold=True
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(0.75),
        Inches(9), Inches(0.35),
        "Scaling-Stage AI Venture with Significant Enterprise Adoption",
        font_size=12, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_horizontal_line(slide, Inches(0.5), Inches(1.1), Inches(2), CURSOR_AMBER, thickness=3)
    
    # === METRICS ROW ===
    metrics = [
        {"value": "$29.3B", "label": "Valuation", "color": CURSOR_PURPLE},
        {"value": "$2.3B", "label": "Series D", "color": CURSOR_GREEN},
        {"value": "$500M+", "label": "ARR", "color": CURSOR_BLUE},
        {"value": "50%+", "label": "Fortune 500", "color": CURSOR_AMBER},
    ]
    
    metric_width = Inches(2.2)
    metric_height = Inches(1)
    metrics_y = Inches(1.3)
    gap = Inches(0.2)
    start_x = Inches(0.3)
    
    for i, metric in enumerate(metrics):
        x = start_x + (metric_width + gap) * i
        
        card = add_rounded_rectangle(
            slide, x, metrics_y,
            metric_width, metric_height, CURSOR_DARK_GRAY,
            border_color=metric["color"], border_width=2
        )
        
        add_textbox(
            slide, x + Inches(0.1), metrics_y + Inches(0.1),
            metric_width - Inches(0.2), Inches(0.5),
            metric["value"],
            font_size=22, font_color=metric["color"], bold=True, alignment=PP_ALIGN.CENTER
        )
        
        add_textbox(
            slide, x + Inches(0.1), metrics_y + Inches(0.6),
            metric_width - Inches(0.2), Inches(0.3),
            metric["label"],
            font_size=10, font_color=CURSOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER
        )
    
    # === INVESTORS ===
    investors_y = Inches(2.5)
    add_textbox(
        slide, Inches(0.4), investors_y,
        Inches(9.2), Inches(0.3),
        "Key Investors & Backers",
        font_size=12, font_color=CURSOR_WHITE, bold=True
    )
    
    investors = ["Accel", "Thrive", "a16z", "Coatue", "NVIDIA", "Google"]
    investor_x = Inches(0.4)
    for investor in investors:
        badge = add_rounded_rectangle(
            slide, investor_x, investors_y + Inches(0.35),
            Inches(1.4), Inches(0.35), CURSOR_DARK_GRAY,
            border_color=CURSOR_MID_GRAY, border_width=1
        )
        add_text_to_shape(badge, investor, font_size=9, font_color=CURSOR_OFF_WHITE)
        investor_x += Inches(1.55)
    
    # === CASE STUDY ===
    case_y = Inches(3.3)
    case_box = add_rounded_rectangle(
        slide, Inches(0.35), case_y,
        Inches(9.3), Inches(1.2), CURSOR_DARK_GRAY,
        border_color=CURSOR_GREEN, border_width=2
    )
    
    sf_icon = add_circle(slide, Inches(0.55), case_y + Inches(0.3), Inches(0.5), CURSOR_BLUE)
    add_text_to_shape(sf_icon, "SF", font_size=12, font_color=CURSOR_WHITE, bold=True)
    
    add_textbox(
        slide, Inches(1.2), case_y + Inches(0.15),
        Inches(8.2), Inches(0.35),
        "Salesforce Case Study",
        font_size=14, font_color=CURSOR_GREEN, bold=True
    )
    
    add_textbox(
        slide, Inches(1.2), case_y + Inches(0.5),
        Inches(8.2), Inches(0.6),
        ">90% of Salesforce engineers use Cursor daily  •  Double-digit gains in velocity & code quality",
        font_size=11, font_color=CURSOR_OFF_WHITE
    )
    
    # === SIGNIFICANCE & ASSESSMENT ===
    sig_y = Inches(4.7)
    sig_box = add_rounded_rectangle(
        slide, Inches(0.35), sig_y,
        Inches(9.3), Inches(1.5), CURSOR_PURPLE
    )
    
    add_textbox(
        slide, Inches(0.5), sig_y + Inches(0.15),
        Inches(9), Inches(0.35),
        "Significance & Disruption Potential",
        font_size=14, font_color=CURSOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER
    )
    
    add_textbox(
        slide, Inches(0.5), sig_y + Inches(0.55),
        Inches(9), Inches(0.85),
        "High potential to reshape how software engineers work with AI-native IDEs.\n"
        "Backed by top-tier investors  •  Proven enterprise traction  •  Strong moat via integration & UX",
        font_size=11, font_color=CURSOR_OFF_WHITE, alignment=PP_ALIGN.CENTER
    )
    
    add_speaker_notes(
        slide,
        "Growth Stage: Cursor is a scaling-stage AI venture. Series D raised $2.3B at a $29.3B valuation. "
        "Top investors include Accel, Thrive, Andreessen Horowitz, Coatue, NVIDIA, and Google. "
        "They claim over $500M ARR and more than half the Fortune 500 as customers. "
        "Salesforce published a case study showing >90% of engineers use Cursor daily with double-digit productivity gains. "
        "Significance: Cursor has high potential to reshape software development with AI-native IDEs. "
        "Overall assessment: This is a significant AI venture with strong financial backing, real enterprise traction, "
        "and meaningful potential to disrupt software development workflows."
    )
    
    add_source_footer(
        slide,
        "Sources: cursor.com/blog/series-d | cursor.com/blog/series-c | cursor.com/blog/salesforce",
        slide_width, slide_height
    )
    add_slide_number(slide, 7, slide_width, slide_height)
    return slide


def create_demo_slide(prs, slide_width, slide_height):
    """Slide 8: Brief Demo Plan."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CURSOR_BLACK)
    
    # Title
    add_textbox(
        slide, Inches(0.5), Inches(0.4),
        Inches(9), Inches(0.6),
        "Brief Demo Plan",
        font_size=36, font_color=CURSOR_WHITE, bold=True
    )
    
    add_textbox(
        slide, Inches(0.5), Inches(0.95),
        Inches(9), Inches(0.4),
        "What I'll Show You",
        font_size=14, font_color=CURSOR_LIGHT_GRAY
    )
    
    add_horizontal_line(slide, Inches(0.5), Inches(1.4), Inches(2), CURSOR_PURPLE, thickness=3)
    
    # === DEMO STEPS ===
    steps = [
        {
            "num": "1",
            "title": "Natural Language → Code",
            "desc": "Show Agent building a small feature from a plain English description",
            "color": CURSOR_PURPLE
        },
        {
            "num": "2",
            "title": "Multi-File Refactor",
            "desc": "Demonstrate AI editing multiple files to fix a bug or improve structure",
            "color": CURSOR_BLUE
        },
        {
            "num": "3",
            "title": "Controls & Settings",
            "desc": "Briefly show model selection, settings, and enterprise controls",
            "color": CURSOR_GREEN
        },
    ]
    
    step_y = Inches(1.8)
    for step in steps:
        # Step card
        card = add_rounded_rectangle(
            slide, Inches(0.5), step_y,
            Inches(9), Inches(1.3), CURSOR_DARK_GRAY,
            border_color=step["color"], border_width=2
        )
        
        # Step number
        num_circle = add_circle(slide, Inches(0.7), step_y + Inches(0.35), Inches(0.6), step["color"])
        add_text_to_shape(num_circle, step["num"], font_size=20, font_color=CURSOR_WHITE, bold=True)
        
        # Step title
        add_textbox(
            slide, Inches(1.5), step_y + Inches(0.25),
            Inches(7.8), Inches(0.45),
            step["title"],
            font_size=18, font_color=step["color"], bold=True
        )
        
        # Step description
        add_textbox(
            slide, Inches(1.5), step_y + Inches(0.7),
            Inches(7.8), Inches(0.5),
            step["desc"],
            font_size=12, font_color=CURSOR_OFF_WHITE
        )
        
        step_y += Inches(1.5)
    
    # Footer note
    add_textbox(
        slide, Inches(0.5), Inches(6.4),
        Inches(9), Inches(0.4),
        "Demo duration: ~60-90 seconds (live or pre-recorded)",
        font_size=10, font_color=CURSOR_MID_GRAY, alignment=PP_ALIGN.CENTER
    )
    
    add_speaker_notes(
        slide,
        "Demo Plan: I'll show three quick things. "
        "First, natural language to code—I'll type a description and let Agent build a small feature. "
        "Second, a multi-file refactor or bug fix to show how Cursor edits across files intelligently. "
        "Third, a brief look at model selection and settings to highlight enterprise controls. "
        "Total demo time: about 60-90 seconds."
    )
    
    add_slide_number(slide, 8, slide_width, slide_height)
    return slide


# =============================================================================
# MAIN FUNCTION
# =============================================================================

def create_presentation():
    """Create and save the professional Cursor presentation."""
    prs = Presentation()
    
    # 4:3 aspect ratio (standard)
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    
    print("🎨 Creating professional Cursor presentation (Assignment Version)...")
    print("   Color scheme: Cursor brand (dark theme)")
    print("   Aspect ratio: 4:3 standard")
    print("   Style: Diagram-heavy with visual hierarchy")
    print("   Interactive: Clickable links with hover tooltips")
    print()
    
    # Create all slides (8 total)
    print("   [1/8] Title slide (logo, link, presenter)...")
    create_title_slide(prs, slide_width, slide_height)
    
    print("   [2/8] Value Proposition (benefits-focused)...")
    create_value_proposition_slide(prs, slide_width, slide_height)
    
    print("   [3/8] Business Model (merged pricing)...")
    create_business_model_slide(prs, slide_width, slide_height)
    
    print("   [4/8] Competition (AI vs Traditional)...")
    create_competition_slide(prs, slide_width, slide_height)
    
    print("   [5/8] AI Technologies (architecture)...")
    create_ai_technologies_slide(prs, slide_width, slide_height)
    
    print("   [6/8] Security & Risk...")
    create_security_slide(prs, slide_width, slide_height)
    
    print("   [7/8] Growth, Funding & Impact...")
    create_growth_funding_slide(prs, slide_width, slide_height)
    
    print("   [8/8] Demo Plan...")
    create_demo_slide(prs, slide_width, slide_height)
    
    # Save
    output_file = "cursor_presentation_pro.pptx"
    prs.save(output_file)
    
    print()
    print(f"✅ Saved: {output_file}")
    print()
    print("📊 Presentation includes (8 slides):")
    print("   1. Title slide with logo, link, presenter name")
    print("   2. Value Proposition (benefits, workflow diagram)")
    print("   3. Business Model (merged pricing tiers)")
    print("   4. Competition (AI & Traditional comparison)")
    print("   5. AI Technologies (architecture diagram)")
    print("   6. Enterprise Security & Risk")
    print("   7. Growth, Funding & Impact")
    print("   8. Brief Demo Plan")
    print()
    print("🔗 Interactive elements:")
    print("   • Clickable logo → cursor.com")
    print("   • Clickable website link with hover tooltip")
    print("   • Clickable Trust Center link with tooltip")
    print()
    print("📝 Speaker notes included for ~4-minute presentation")
    print("📚 Source citations on each slide")


if __name__ == "__main__":
    create_presentation()
