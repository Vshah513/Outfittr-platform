#!/usr/bin/env python3
"""
Create a professional 7-slide Cursor presentation using python-pptx.
Includes speaker notes, source footers, and a pricing table.
Run: python3 create_cursor_presentation.py
Output: cursor_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


# Color palette (professional tech)
PRIMARY_BLUE = RGBColor(0, 122, 255)  # Cursor-inspired blue
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)
ACCENT_ORANGE = RGBColor(255, 149, 0)


def add_source_footer(slide, source_text):
    """Add a small source footer at the bottom of the slide."""
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = source_text
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.font.italic = True


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_title_and_subtitle(slide, title, subtitle=""):
    """Add title and optional subtitle to a slide."""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    
    # Subtitle (if provided)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(1))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY


def add_bullets(slide, bullets, start_y=1.8):
    """Add bullet points to a slide."""
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(8.4), Inches(4.5))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)
        p.level = 0


def add_pricing_table(slide):
    """Add a 5x5 pricing table to slide 4."""
    # Table: Plans x Features
    rows, cols = 5, 5
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(4.2)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Define data
    headers = ["Feature", "Hobby", "Pro", "Pro+", "Ultra"]
    data = [
        headers,
        ["Monthly Cost", "Free", "$20", "$60", "$200"],
        ["Agent Requests", "Limited", "Extended", "Extended", "Extended"],
        ["Tab Completions", "Limited", "Unlimited", "Unlimited", "Unlimited"],
        ["Model Usage", "Basic", "1x", "3x", "20x"],
    ]
    
    # Populate table
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text
            
            # Style header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_BLUE
                text_frame = cell.text_frame
                text_frame.paragraphs[0].font.size = Pt(12)
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].font.color.rgb = WHITE
            else:
                # Alternate row colors
                cell.fill.solid()
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                else:
                    cell.fill.fore_color.rgb = WHITE
                
                text_frame = cell.text_frame
                text_frame.paragraphs[0].font.size = Pt(11)
                text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
            
            text_frame = cell.text_frame
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def create_presentation():
    """Create and save the Cursor presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    add_title_and_subtitle(
        slide,
        "Cursor: AI Coding IDE",
        "4-minute overview: What it is, how it works, pricing, enterprise posture, traction"
    )
    
    add_speaker_notes(
        slide,
        "Welcome. We're covering Cursor, the AI coding IDE. In 4 minutes, you'll see what it does, "
        "how the AI powers it, pricing options, enterprise security, and why it's growing fast. Let's start."
    )
    
    # ========== SLIDE 2: What Cursor is (Value Prop) ==========
    slide = prs.slides.add_slide(blank_layout)
    add_title_and_subtitle(slide, "What Cursor is: AI Editor + Coding Agent")
    
    bullets = [
        "Describe what you want in natural language; Cursor writes the code",
        "Tab: Intelligent autocompletion that learns from your accept/reject feedback",
        "Agent: Completes complex tasks independently—edits code, runs terminal commands, debugs",
        "Access Agent via Cmd/Ctrl+I sidepane for hands-off coding assistance"
    ]
    add_bullets(slide, bullets)
    
    add_speaker_notes(
        slide,
        "Cursor is an AI editor. You describe what you want, and it writes code. Two main features: "
        "Tab is smart autocompletion—it gets smarter as you accept or reject its suggestions. "
        "Agent is a true coding assistant that can complete tasks by itself—it edits files, runs commands, fixes bugs. "
        "You trigger it with Cmd or Ctrl+I in the sidebar. Think of it as a pair programmer that works independently."
    )
    
    add_source_footer(
        slide,
        "Sources: https://cursor.com/docs | https://cursor.com/docs/agent/overview | https://cursor.com/docs/tab/overview"
    )
    
    # ========== SLIDE 3: How the AI Works ==========
    slide = prs.slides.add_slide(blank_layout)
    add_title_and_subtitle(slide, "AI Backbone: Multi-Provider Models")
    
    bullets = [
        "Supports frontier models from OpenAI, Anthropic, Google, and others",
        "Cursor compares capabilities, context windows, and pricing across providers",
        "Agent can work with any supported model; pick based on your task and budget",
        "Auto mode: Cursor intelligently routes requests to optimal model based on capacity/cost"
    ]
    add_bullets(slide, bullets)
    
    add_speaker_notes(
        slide,
        "Cursor doesn't build its own model. Instead, it integrates frontier models from top providers: "
        "OpenAI, Anthropic, Google. You can see a comparison of their capabilities and costs. "
        "The Agent works with any of these. There's also an Auto mode that picks the best model automatically based on what's available and what's cheapest. "
        "This flexibility means you get state-of-the-art AI without vendor lock-in."
    )
    
    add_source_footer(slide, "Source: https://cursor.com/docs/models")
    
    # ========== SLIDE 4: Pricing (with table) ==========
    slide = prs.slides.add_slide(blank_layout)
    add_title_and_subtitle(slide, "Pricing: Individual Plans")
    
    add_pricing_table(slide)
    
    # Add note below table
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "• Pricing shifted to usage-based model: Pro includes $20/mo of usage. Auto option enables unlimited usage by rotating models."
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    
    add_speaker_notes(
        slide,
        "Cursor offers four individual plans. Hobby is free—good for trying it out. "
        "Pro is 20 a month and gives you extended Agent requests and unlimited Tab completions. "
        "Pro Plus adds 3x usage multiplier on OpenAI and Anthropic models for 60 a month. "
        "Ultra, at 200 a month, gives you 20x multiplier and priority access to new features. "
        "A recent change: Cursor moved from request-based to usage-based pricing. Pro includes 20 dollars of usage; you can go over that with Auto mode, which spreads requests across models to stay efficient."
    )
    
    add_source_footer(
        slide,
        "Sources: https://cursor.com/pricing | https://cursor.com/blog/june-2025-pricing"
    )
    
    # ========== SLIDE 5: Teams & Enterprise ==========
    slide = prs.slides.add_slide(blank_layout)
    add_title_and_subtitle(slide, "Teams & Enterprise Packaging")
    
    bullets = [
        "Teams ($40/user/mo): Shared chats, centralized billing, RBAC, SAML/OIDC SSO, usage analytics",
        "Enterprise (custom): Pooled usage, invoice/PO billing, SCIM, AI audit logs, granular model controls, priority support",
        "Both include on-demand usage beyond monthly seat allowance",
        "Enterprise gets dedicated admin and security compliance controls"
    ]
    add_bullets(slide, bullets)
    
    add_speaker_notes(
        slide,
        "For teams, Cursor offers Teams at 40 per user per month. You get shared chats and commands, "
        "centralized billing, role-based access, and single sign-on via SAML or OIDC. "
        "Enterprise pricing is custom. You get pooled usage, invoice billing, full audit logs of AI code, "
        "granular controls over which models your team can use, and priority support. "
        "Both plans include overage allowances if you exceed your monthly seat usage."
    )
    
    add_source_footer(
        slide,
        "Sources: https://cursor.com/pricing | https://cursor.com/docs/account/teams/pricing | https://cursor.com/docs/account/teams/sso"
    )
    
    # ========== SLIDE 6: Enterprise Security Posture ==========
    slide = prs.slides.add_slide(blank_layout)
    add_title_and_subtitle(slide, "Enterprise Security & Compliance")
    
    bullets = [
        "SOC 2 Type II certified; third-party pen tests at least annually (details in Trust Center)",
        "AES-256 encryption at rest; TLS 1.2+ in transit",
        "Centralized security controls, GDPR/CCPA compliance references",
        "Trust Center provides audit resources on request: https://trust.cursor.com"
    ]
    add_bullets(slide, bullets)
    
    add_speaker_notes(
        slide,
        "On security: Cursor is SOC 2 Type II certified. They commit to annual third-party penetration testing; "
        "details available through their Trust Center. Data is encrypted at rest using AES-256 and in transit with TLS 1.2 or higher. "
        "Enterprise customers get centralized security controls and compliance references for GDPR and CCPA. "
        "If you need audit reports or more details, the Trust Center at trust.cursor.com has them available on request."
    )
    
    add_source_footer(
        slide,
        "Sources: https://cursor.com/security | https://cursor.com/enterprise | https://trust.cursor.com"
    )
    
    # ========== SLIDE 7: Traction & Funding + Closing ==========
    slide = prs.slides.add_slide(blank_layout)
    add_title_and_subtitle(slide, "Traction & Funding: Rapid Growth")
    
    bullets = [
        "Series D (2025): $2.3B raised at $29.3B valuation; investors: Accel, Thrive, A16z, Coatue, NVIDIA, Google",
        "Series C: $900M at $9.9B valuation; claimed >$500M ARR and >50% of Fortune 500 daily users",
        "Salesforce case study: >90% of Salesforce engineers use Cursor; reported double-digit gains in velocity & code quality",
        "💡 Bottom line: Cursor is scaling fast. Evaluate fit by workflow fit, security needs, and cost model."
    ]
    add_bullets(slide, bullets)
    
    add_speaker_notes(
        slide,
        "Cursor's growth is impressive. In series D, they raised two point three billion at a valuation of 29 point three billion. "
        "Top-tier investors are backing them: Accel, Thrive, Andreessen Horowitz, Coatue, NVIDIA, and Google. "
        "In their Series C, they claimed over 500 million in annual revenue and more than half the Fortune 500 as daily users. "
        "Salesforce published a case study showing over 90 percent of their engineers use Cursor daily and saw double-digit improvements in code velocity and quality. "
        "Bottom line: Cursor is scaling fast. If you're evaluating it, focus on whether it fits your workflow, "
        "whether the security and compliance story works for you, and whether the pricing model makes sense for your usage."
    )
    
    add_source_footer(
        slide,
        "Sources: https://cursor.com/blog/series-d | https://cursor.com/blog/series-c | https://cursor.com/blog/salesforce"
    )
    
    # Save
    prs.save("cursor_presentation.pptx")
    print("✅ Saved: cursor_presentation.pptx")
    print("📊 7 slides with speaker notes, sources, and pricing table")


if __name__ == "__main__":
    create_presentation()
